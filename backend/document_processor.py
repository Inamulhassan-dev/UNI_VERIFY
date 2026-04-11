"""
Universal Document Processor for UNI-VERIFY
Handles: PDF, DOCX, TXT, PPTX → compressed .txt.gz storage
"""

import os
import gzip
import uuid
import re
from typing import Optional, Tuple

import fitz  # PyMuPDF (for PDFs)

# Try importing optional document libraries
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("[WARN] python-docx not installed. DOCX support disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("[WARN] python-pptx not installed. PPTX support disabled.")


# Supported file extensions
SUPPORTED_FORMATS = {'.pdf', '.docx', '.doc', '.txt', '.pptx'}

# Data warehouse directory
WAREHOUSE_DIR = os.path.join(os.path.dirname(__file__), "data_warehouse")
os.makedirs(WAREHOUSE_DIR, exist_ok=True)


# ============ TEXT EXTRACTION ============

def extract_text_from_file(file_path: str) -> Optional[str]:
    """
    Extract text from any supported file format.
    
    Args:
        file_path: Path to the uploaded file
        
    Returns:
        Extracted and cleaned text, or None if extraction fails
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return _extract_from_pdf(file_path)
    elif ext in ('.docx', '.doc'):
        return _extract_from_docx(file_path)
    elif ext == '.txt':
        return _extract_from_txt(file_path)
    elif ext == '.pptx':
        return _extract_from_pptx(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return None


def _extract_from_pdf(file_path: str) -> Optional[str]:
    """Extract text from PDF using PyMuPDF."""
    try:
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return _clean_text('\n'.join(text_parts))
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return None


def _extract_from_docx(file_path: str) -> Optional[str]:
    """Extract text from DOCX (Word) files."""
    if not DOCX_AVAILABLE:
        print("Cannot process DOCX: python-docx not installed")
        return None
    try:
        doc = DocxDocument(file_path)
        text_parts = []
        
        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return _clean_text('\n'.join(text_parts))
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
        return None


def _extract_from_txt(file_path: str) -> Optional[str]:
    """Extract text from plain text files."""
    try:
        # Try multiple encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252', 'ascii']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                return _clean_text(text)
            except UnicodeDecodeError:
                continue
        print("Could not decode text file with any known encoding")
        return None
    except Exception as e:
        print(f"Error extracting TXT: {e}")
        return None


def _extract_from_pptx(file_path: str) -> Optional[str]:
    """Extract text from PowerPoint files."""
    if not PPTX_AVAILABLE:
        print("Cannot process PPTX: python-pptx not installed")
        return None
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
                # Extract from tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            text_parts.append(row_text)
        
        return _clean_text('\n'.join(text_parts))
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        return None


# ============ TEXT CLEANING ============

def _clean_text(text: str) -> str:
    """Clean and normalize extracted text."""
    if not text:
        return ""
    
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Remove special characters but keep basic punctuation
    text = re.sub(r'[^\w\s.,!?;:\-\'\"()\[\]]', '', text)
    
    # Remove very short fragments (noise)
    words = text.split()
    if len(words) < 5:
        return ""
    
    return text.strip()


# ============ COMPRESSED STORAGE ============

def save_compressed_text(text: str, project_id: str = None) -> Tuple[str, int]:
    """
    Save extracted text as a compressed .txt.gz file.
    
    Args:
        text: Cleaned text content
        project_id: Optional ID for the filename
        
    Returns:
        Tuple of (file_path, compressed_size_bytes)
    """
    if project_id is None:
        project_id = str(uuid.uuid4())[:12]
    
    filename = f"proj_{project_id}.txt.gz"
    file_path = os.path.join(WAREHOUSE_DIR, filename)
    
    # Compress and save
    text_bytes = text.encode('utf-8')
    compressed = gzip.compress(text_bytes, compresslevel=9)  # Max compression
    
    with open(file_path, 'wb') as f:
        f.write(compressed)
    
    original_size = len(text_bytes)
    compressed_size = len(compressed)
    ratio = (1 - compressed_size / original_size) * 100 if original_size > 0 else 0
    
    print(f"[SAVED] {filename} | Original: {original_size:,}B -> Compressed: {compressed_size:,}B ({ratio:.1f}% smaller)")
    
    return file_path, compressed_size


def load_compressed_text(file_path: str) -> Optional[str]:
    """
    Load text from a compressed .txt.gz file.
    
    Args:
        file_path: Path to the .txt.gz file
        
    Returns:
        Decompressed text content
    """
    try:
        with gzip.open(file_path, 'rb') as f:
            return f.read().decode('utf-8')
    except Exception as e:
        print(f"Error loading compressed text: {e}")
        return None


def delete_compressed_text(file_path: str) -> bool:
    """Delete a compressed text file from the warehouse."""
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            return True
        return False
    except Exception as e:
        print(f"Error deleting file: {e}")
        return False


# ============ UTILITIES ============

def get_file_format(filename: str) -> str:
    """Get the file extension/format."""
    return os.path.splitext(filename)[1].lower().lstrip('.')


def is_supported_format(filename: str) -> bool:
    """Check if a file format is supported."""
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_FORMATS


def get_warehouse_stats() -> dict:
    """Get statistics about the data warehouse storage."""
    total_files = 0
    total_size = 0
    
    if os.path.exists(WAREHOUSE_DIR):
        for f in os.listdir(WAREHOUSE_DIR):
            if f.endswith('.txt.gz'):
                total_files += 1
                total_size += os.path.getsize(os.path.join(WAREHOUSE_DIR, f))
    
    return {
        "total_files": total_files,
        "total_size_bytes": total_size,
        "total_size_mb": round(total_size / (1024 * 1024), 2) if total_size > 0 else 0,
        "warehouse_path": WAREHOUSE_DIR
    }


def validate_document_content(text: str, min_words: int = 50) -> Tuple[bool, str]:
    """Validate that the document has enough content for analysis."""
    if not text:
        return False, "Could not extract text from the document"
    
    word_count = len(text.split())
    
    if word_count < min_words:
        return False, f"Document has insufficient content ({word_count} words). Minimum {min_words} words required."
    
    return True, f"Document validated ({word_count} words)"


def extract_title_and_abstract(text: str) -> dict:
    """
    Try to extract title and abstract from the text.
    Uses common patterns in academic documents.
    """
    result = {
        "title": "",
        "abstract": ""
    }
    
    # Try to find title (usually first significant line)
    lines = text.split('.')
    if lines:
        # First sentence that's reasonably long could be title
        for line in lines[:5]:
            line = line.strip()
            if 10 < len(line) < 200:
                result["title"] = line
                break
    
    # Try to find abstract section
    abstract_pattern = r'abstract[:\s]*(.{100,1000}?)(?:introduction|keywords|1\.|chapter|$)'
    abstract_match = re.search(abstract_pattern, text.lower())
    
    if abstract_match:
        result["abstract"] = abstract_match.group(1).strip()
    else:
        # If no abstract found, use first 500 chars as summary
        result["abstract"] = text[:500] if len(text) > 500 else text
    
    return result
